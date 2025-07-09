
from pathlib import Path
import os
import json
import subprocess
import whisper
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation
import pandas as pd
import torch
from transformers import BlipProcessor, BlipForConditionalGeneration
import cv2
import docx2txt


VIDEO_EXTENSIONS = {'.mp4', '.mov', '.avi', '.mkv', '.flv', '.wmv', '.webm', '.m4v'}
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp', '.gif', '.svg'}
PDF_EXTENSIONS = {'.pdf'}
SLIDE_EXTENSIONS = {'.pptx', '.ppt', '.odp', '.key'}
SHEET_EXTENSIONS = {'.xlsx', '.xls', '.ods', '.csv', '.tsv'}
DOC_EXTENSIONS = {'.docx', '.doc', '.odt'}
TEXT_EXTENSIONS = {'.txt', '.md', '.rtf', '.log'}

whisper_model = whisper.load_model("base")
blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")

def extract_scene_descriptions(video_path, interval_seconds=5):
    video = cv2.VideoCapture(str(video_path))
    fps = video.get(cv2.CAP_PROP_FPS)
    interval_frames = int(fps * interval_seconds)
    frame_count = 0
    descriptions = []

    while True:
        success, frame = video.read()
        if not success:
            break
        if frame_count % interval_frames == 0:
            image = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
            inputs = blip_processor(image, return_tensors="pt")
            with torch.no_grad():
                out = blip_model.generate(**inputs)
            caption = blip_processor.decode(out[0], skip_special_tokens=True)
            descriptions.append(caption)
        frame_count += 1

    video.release()
    return descriptions

def extract_video_metadata(file_path):
    cmd = [
        "ffprobe", "-v", "error", "-select_streams", "v:0",
        "-show_entries", "stream=width,height,duration",
        "-of", "default=noprint_wrappers=1", str(file_path)
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    meta = dict(line.split('=') for line in result.stdout.strip().split('\n') if '=' in line)

    try:
        transcription = whisper_model.transcribe(str(file_path))["text"]
    except Exception:
        transcription = ""

    try:
        scene_descriptions = extract_scene_descriptions(file_path)
    except Exception:
        scene_descriptions = []

    return {
        "documentId": file_path.stem,
        "title": file_path.name,
        "content": transcription,
        "metadata": {
            "type": "video",
            "source": str(file_path),
            "duration": meta.get("duration", ""),
            "resolution": f"{meta.get('width')}x{meta.get('height')}",
            "transcript": transcription,
            "scene_descriptions": scene_descriptions
        }
    }

def extract_image_metadata(file_path):
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
    except Exception:
        text = ""
    return {
        "documentId": file_path.stem,
        "title": file_path.name,
        "content": text,
        "metadata": {
            "type": "image",
            "source": str(file_path),
        }
    }

def extract_pdf_metadata(file_path):
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
    except Exception:
        text = ""
    return {
        "documentId": file_path.stem,
        "title": file_path.name,
        "content": text,
        "metadata": {
            "type": "pdf",
            "source": str(file_path)
        }
    }

def extract_pptx_metadata(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
    except Exception:
        text = ""
    return {
        "documentId": file_path.stem,
        "title": file_path.name,
        "content": text,
        "metadata": {
            "type": "slide",
            "source": str(file_path)
        }
    }

def extract_spreadsheet_metadata(file_path):
    text = ""
    try:
        xls = pd.read_excel(file_path, sheet_name=None)
        for name, df in xls.items():
            text += f"Sheet: {name}\n{df.to_string()}\n\n"
    except Exception:
        text = ""
    return {
        "documentId": file_path.stem,
        "title": file_path.name,
        "content": text,
        "metadata": {
            "type": "spreadsheet",
            "source": str(file_path)
        }
    }

def extract_text_file_metadata(file_path):
    text = ""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
    except Exception:
        text = ""
    return {
        "documentId": file_path.stem,
        "title": file_path.name,
        "content": text,
        "metadata": {
            "type": "text",
            "source": str(file_path)
        }
    }

def extract_doc_file_metadata(file_path):
    text = ""
    try:
        if file_path.suffix.lower() == ".docx":
            text = docx2txt.process(str(file_path))
        else:
            text = f"(Unsupported DOC format: {file_path.suffix})"
    except Exception:
        text = ""
    return {
        "documentId": file_path.stem,
        "title": file_path.name,
        "content": text,
        "metadata": {
            "type": "document",
            "source": str(file_path)
        }
    }

def save_json(data, output_dir):
    output_dir.mkdir(parents=True, exist_ok=True)
    json_path = output_dir / f"{data['documentId']}.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)

def process_all_files(root_dir, output_dir):
    root_path = Path(root_dir)
    for file_path in root_path.rglob("*.*"):
        if not file_path.is_file():
            continue

        suffix = file_path.suffix.lower()
        try:
            if suffix in VIDEO_EXTENSIONS:
                data = extract_video_metadata(file_path)
            elif suffix in IMAGE_EXTENSIONS:
                data = extract_image_metadata(file_path)
            elif suffix in PDF_EXTENSIONS:
                data = extract_pdf_metadata(file_path)
            elif suffix in SLIDE_EXTENSIONS:
                data = extract_pptx_metadata(file_path)
            elif suffix in SHEET_EXTENSIONS:
                data = extract_spreadsheet_metadata(file_path)
            elif suffix in TEXT_EXTENSIONS:
                data = extract_text_file_metadata(file_path)
            elif suffix in DOC_EXTENSIONS:
                data = extract_doc_file_metadata(file_path)
            
            else:
                continue
            save_json(data, output_dir)
            print(f"Processed: {file_path}")
        except Exception as e:
            print(f"Failed to process {file_path}: {e}")


root_folder = "FOLDER_PATH"
output_folder = Path(root_folder) / "metadata-jsons"
process_all_files(root_folder, output_folder)
